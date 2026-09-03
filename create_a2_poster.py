from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

BASE = r'd:\GraduateDesign\Test\dcase2025_task6_baseline\ServerCodes\TestVersion\d25_t6'
OUT = BASE + r'\A2_Audio_Text_Retrieval_Poster_v5.pptx'
IMG = BASE + '\\images\\'

prs = Presentation()
prs.slide_width = Inches(16.54)
prs.slide_height = Inches(23.39)
slide = prs.slides.add_slide(prs.slide_layouts[6])

NAVY=RGBColor(18,38,70); TEAL=RGBColor(23,129,146); ORANGE=RGBColor(233,137,40)
BG=RGBColor(245,247,250); WHITE=RGBColor(255,255,255); DARK=RGBColor(35,40,45); GREY=RGBColor(95,103,112)
GREEN=RGBColor(56,151,95); PURPLE=RGBColor(118,91,180); RED=RGBColor(188,74,74); LINE=RGBColor(218,226,233)
LBLUE=RGBColor(240,247,252); LPURPLE=RGBColor(247,242,251); LORANGE=RGBColor(255,248,236); LGREEN=RGBColor(239,250,242)


def rect(x,y,w,h,fill=WHITE,line=LINE,round=True,linew=1):
    s=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if round else MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb=fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb=line; s.line.width=Pt(linew)
    return s


def text(t,x,y,w,h,size=14,b=False,c=DARK,center=False):
    tb=slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf=tb.text_frame; tf.clear(); tf.word_wrap=True
    tf.margin_left=Inches(0.05); tf.margin_right=Inches(0.05); tf.margin_top=Inches(0.03); tf.margin_bottom=Inches(0.03)
    p=tf.paragraphs[0]; p.text=t; p.font.name='Arial'; p.font.size=Pt(size); p.font.bold=b; p.font.color.rgb=c
    if center: p.alignment=PP_ALIGN.CENTER
    return tb


def bullet(items,x,y,w,h,size=11.6):
    tb=slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf=tb.text_frame; tf.clear(); tf.word_wrap=True
    tf.margin_left=Inches(0.08); tf.margin_right=Inches(0.04); tf.margin_top=Inches(0.03)
    for i,it in enumerate(items):
        p=tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text='• '+it; p.font.name='Arial'; p.font.size=Pt(size); p.font.color.rgb=DARK; p.space_after=Pt(4)
    return tb


def titlebar(t,x,y,w,col):
    rect(x,y,w,0.42,col,None)
    text(t,x+0.08,y+0.06,w-0.16,0.27,14.2,True,WHITE)


def metric(v,l,x,y,w=1.42,col=TEAL):
    rect(x,y,w,1.02,WHITE,LINE)
    text(v,x,y+0.12,w,0.34,23,True,col,True)
    text(l,x+0.05,y+0.60,w-0.1,0.24,10.0,False,GREY,True)


def node(t,x,y,w,h,col,bg):
    rect(x,y,w,h,bg,col,True,1.35)
    text(t,x+0.03,y+0.08,w-0.06,h-0.12,10.9,True,col,True)


def arrow(x1,y1,x2,y2,col=GREY,w=1.8):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = col
    line.line.width = Pt(w)
    try:
        line.line.end_arrowhead = True
    except Exception:
        pass
    return line

# background and header
rect(0,0,16.54,23.39,BG,None,False)
rect(0,0,16.54,2.08,NAVY,None,False)
text('Design and Evaluation of a Multi-Encoder\nAudio-Text Retrieval System',0.75,0.24,15.0,0.92,29,True,WHITE,True)
text('Combining PaSST and CLAP for retrieval-based audio-text matching',1.3,1.23,13.9,0.28,14.5,False,RGBColor(220,231,242),True)
text('Anonymous Copy',1.3,1.57,13.9,0.2,11.8,False,RGBColor(196,211,228),True)

# left column
titlebar('1. Background & Objective',0.45,2.42,3.9,NAVY)
rect(0.45,2.88,3.9,3.0)
text('Motivation',0.66,3.08,3.2,0.22,14.7,True,NAVY)
text('Traditional audio classification relies on fixed labels, which limits flexibility when users need natural language descriptions of diverse acoustic scenes.',0.66,3.42,3.45,0.66,12.3)
text('Objective',0.66,4.30,3.2,0.22,14.7,True,NAVY)
text('Design and evaluate a retrieval-based system that combines complementary audio encoders and retrieves relevant text descriptions from a large candidate library.',0.66,4.64,3.45,0.70,12.3)

titlebar('2. Data & Candidate Library',0.45,6.15,3.9,NAVY)
metric('4,981','Clotho clips',0.56,6.67,1.07)
metric('24,905','Clotho captions',1.74,6.67,1.07)
metric('80,045','candidate descriptions',2.92,6.67,1.07,ORANGE)
rect(0.45,7.90,3.9,1.55)
bullet(['Clotho was used for evaluation.','AudioCaps expanded the candidate description library.','Filtering and deduplication were applied before inference.'],0.60,8.14,3.55,1.14,11.4)

titlebar('3. Experimental Development',0.45,9.72,3.9,NAVY)
metric('23','iterations',0.56,10.22,1.07,PURPLE)
metric('11','retained',1.74,10.22,1.07,GREEN)
metric('12','discarded',2.92,10.22,1.07,RED)
rect(0.45,11.42,3.9,1.90)
bullet(['Built and refined the retrieval system through 23 experimental iterations.','Compared baselines, fusion mechanisms, projection heads, and training settings.','Retained only the configurations that improved stability or retrieval quality.'],0.63,11.68,3.52,1.36,11.3)

titlebar('4. Contributions',0.45,13.58,3.9,NAVY)
rect(0.45,14.04,3.9,2.55)
bullet(['Developed a complete audio-text retrieval system for matching audio with natural-language descriptions.','Built an 80,045-description candidate library for large-scale retrieval.','Improved retrieval performance to 0.325 mAP@10 through systematic refinement during development.','Deployed the final system as a practical real-time web application.'],0.62,14.30,3.54,1.98,11.1)

# center top method
titlebar('5. Proposed Method',4.60,2.42,6.85,TEAL)
rect(4.60,2.88,6.85,7.55)
text('System Architecture',4.88,3.08,2.6,0.24,15.7,True,NAVY)
text('Audio processing branch',4.95,3.48,2.25,0.20,12.5,True,TEAL)
text('Text retrieval branch',8.86,3.48,2.25,0.20,12.5,True,PURPLE)
node('Input Audio',5.32,3.86,1.55,0.54,TEAL,LBLUE)
node('PaSST',4.86,4.86,1.38,0.54,GREEN,LGREEN)
node('CLAP',5.98,4.86,1.38,0.54,GREEN,LGREEN)
node('Attention Fusion',5.28,6.00,1.70,0.60,ORANGE,LORANGE)
node('Audio Embedding',5.30,7.18,1.66,0.58,TEAL,LBLUE)
node('Candidate Captions',8.98,3.86,1.72,0.54,PURPLE,LPURPLE)
node('RoBERTa-large',8.98,4.98,1.72,0.54,PURPLE,LPURPLE)
node('Text Embedding',8.98,6.12,1.72,0.60,PURPLE,LPURPLE)
node('Library Matrix',8.98,7.30,1.72,0.58,PURPLE,LPURPLE)
node('Cosine Similarity',7.12,8.00,2.02,0.62,ORANGE,LORANGE)
node('Ranked Text Descriptions',6.45,9.02,3.32,0.64,GREEN,LGREEN)
arrow(6.08,4.40,5.55,4.86,TEAL)
arrow(6.10,4.40,6.66,4.86,TEAL)
arrow(5.55,5.40,6.05,6.00,GREEN)
arrow(6.66,5.40,6.18,6.00,GREEN)
arrow(6.13,6.60,6.13,7.18,ORANGE)
arrow(9.84,4.40,9.84,4.98,PURPLE)
arrow(9.84,5.52,9.84,6.12,PURPLE)
arrow(9.84,6.72,9.84,7.30,PURPLE)
arrow(6.96,7.76,7.32,8.00,TEAL)
arrow(9.00,7.88,8.92,8.00,PURPLE)
arrow(8.13,8.62,8.13,9.02,ORANGE)
text('PaSST captures fine-grained time-frequency patterns, while CLAP contributes semantic audio-language alignment. Their complementary embeddings are fused and matched against a RoBERTa-large encoded candidate library.',4.95,9.82,6.15,0.45,11.7)

# center deployment
titlebar('6. Deployment Example',4.60,10.72,6.85,TEAL)
rect(4.60,11.18,6.85,4.65)
slide.shapes.add_picture(IMG+'RainAndBark.png', Inches(4.98), Inches(11.55), width=Inches(6.10))
text('Representative deployment screenshot: the web application accepts audio input and returns ranked natural-language descriptions with similarity scores.',4.95,15.10,6.10,0.42,11.5)

# right results
titlebar('7. Quantitative Results',11.75,2.42,4.34,ORANGE)
rect(11.75,2.88,4.34,3.45)
text('Method',12.02,3.16,1.25,0.20,11.1,True,NAVY)
text('mAP@10   R@1   R@10',13.16,3.16,2.45,0.20,11.1,True,NAVY)
rows=[('PaSST-only','0.290   0.181   0.588',False),('CLAP-only','0.271   0.161   0.560',False),('Simple fusion','0.283   0.174   0.562',False),('Final system','0.325   0.203   0.623',True)]
for i,(a,b,hi) in enumerate(rows):
    y=3.58+i*0.40
    text(a,12.02,y,1.45,0.20,10.7,hi,ORANGE if hi else DARK)
    text(b,13.16,y,2.45,0.20,10.7,hi,ORANGE if hi else DARK)
metric('+12.1%','over PaSST-only',12.12,5.22,1.68,ORANGE)
metric('+20.0%','over CLAP-only',13.98,5.22,1.68,ORANGE)

# right training curves
titlebar('8. Training Curves',11.75,6.62,4.34,ORANGE)
rect(11.75,7.08,4.34,5.95)
slide.shapes.add_picture(IMG+'ValmAP.png', Inches(11.98), Inches(7.35), width=Inches(3.85))
slide.shapes.add_picture(IMG+'ValR@10.png', Inches(11.98), Inches(9.85), width=Inches(3.85))
text('Validation mAP and Recall@10 curves show stable convergence after iterative tuning and training-stability improvements.',12.02,12.32,3.85,0.42,11.2)

# right ablation
titlebar('9. Ablation & Key Findings',11.75,13.28,4.34,ORANGE)
rect(11.75,13.74,4.34,2.20)
bullet(['Attention-based text pooling produced the largest individual gain.','Learned fusion outperformed simple concatenation.','The final system reached 0.325 mAP@10 after 23 iterations.'],11.95,13.98,3.92,1.45,11.2)

# bottom content revised
titlebar('10. Limitations, Future Work & Project Outcome',0.45,16.60,15.65,NAVY)
rect(0.45,17.05,15.65,2.65)
text('Limitations',0.75,17.34,2.2,0.20,14.0,True,NAVY)
bullet(['Limited by Clotho training scale','Rare events depend on library coverage','English-only descriptions'],0.75,17.72,4.0,1.08,11.2)
text('Future Work',5.80,17.34,2.2,0.20,14.0,True,NAVY)
bullet(['WavCaps pretraining','Diverse encoder ensembles','Stabilised cross-modal attention','Multilingual retrieval'],5.80,17.72,4.05,1.08,11.2)
text('Project Outcome',10.72,17.34,2.7,0.20,14.0,True,NAVY)
text('A complete and deployable audio-text retrieval system was built, evaluated, and refined through systematic experimentation.',10.72,17.74,4.55,0.52,11.2)

rect(0.45,20.02,15.65,1.30,RGBColor(234,243,248),LINE)
text('Key Takeaway',0.80,20.30,2.1,0.22,14.2,True,NAVY)
text('The project demonstrates that attention-based fusion of PaSST and CLAP improves audio-text retrieval and enables practical real-time deployment on standard hardware.',2.30,20.24,13.10,0.38,12.5,True,ORANGE,True)

rect(0,21.52,16.54,1.87,RGBColor(249,250,252),None,False)
rect(0,22.55,16.54,0.84,NAVY,None,False)
text('Design and Evaluation of a Multi-Encoder Audio-Text Retrieval System',0.8,22.76,14.9,0.24,12.5,False,RGBColor(230,240,250),True)

prs.save(OUT)
print(OUT)
